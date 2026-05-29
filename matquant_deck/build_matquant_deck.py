from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "to_human"
OUT = OUT_DIR / "MatQuant_4page_intro.pptx"
TEMPLATE = ROOT / "template.pptx"

NAVY = RGBColor(8, 56, 118)
BLUE = RGBColor(91, 155, 213)
LIGHT_BLUE = RGBColor(222, 236, 249)
ORANGE = RGBColor(237, 125, 49)
LIGHT_ORANGE = RGBColor(255, 235, 220)
GREEN = RGBColor(112, 173, 71)
LIGHT_GREEN = RGBColor(229, 244, 220)
YELLOW = RGBColor(243, 226, 127)
RED = RGBColor(192, 80, 77)
PURPLE = RGBColor(112, 48, 160)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(243, 245, 248)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def remove_shape(shape) -> None:
    el = shape._element  # pylint: disable=protected-access
    el.getparent().remove(el)


def disable_bullet(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()  # pylint: disable=protected-access
    for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buNone"):
        for child in list(p_pr.findall(qn(tag))):
            p_pr.remove(child)
    p_pr.append(OxmlElement("a:buNone"))


def set_runs(paragraph, font_size=14, color=BLACK, bold=False, italic=False) -> None:
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    for run in paragraph.runs:
        if run.font is None:
            continue
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 14,
    color=BLACK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_runs(p, font_size, color, bold, italic)
    return box


def add_slide(prs: Presentation, title: str, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[17])
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx != 0:
            remove_shape(shape)
    title_shape = slide.shapes.title
    title_shape.text = title
    tf = title_shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    for p in tf.paragraphs:
        disable_bullet(p)
        set_runs(p, 29, BLACK, False)
    add_textbox(slide, str(slide_no), 15.42, 8.28, 0.42, 0.30, font_size=10, color=WHITE)
    return slide


def add_goal(slide, text: str, y: float = 1.72) -> None:
    add_textbox(slide, text, 1.12, y, 9.6, 0.35, font_size=15, color=NAVY, bold=True)
    for i in range(8):
        dash = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1.12 + i * 0.82),
            Inches(y + 0.43),
            Inches(0.58),
            Inches(0.07),
        )
        dash.fill.solid()
        dash.fill.fore_color.rgb = BLUE
        dash.line.fill.background()


def add_card(slide, x, y, w, h, title, bullets, accent=NAVY, fill=WHITE, title_size=14, body_size=11.2):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.4)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.10)
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, title_size, accent, True)
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.space_before = Pt(5)
        p.level = 0
        set_runs(p, body_size, BLACK, False)
    return shape


def add_label(slide, text, x, y, w, h, fill, font_size=12, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, font_size, color, True)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.4):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_callout(slide, text, x=2.1, y=7.08, w=7.8, h=0.46):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = YELLOW
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, 13, NAVY, True)


def add_source(slide, text="Source: Nair et al., Matryoshka Quantization, ICML/PMLR 2025; arXiv:2502.06786"):
    add_textbox(slide, text, 1.12, 8.02, 9.7, 0.30, font_size=10, color=GRAY)


def build_slide_1(prs: Presentation) -> None:
    slide = add_slide(prs, "MatQuant: Problem Setting and Research Gap", 1)
    add_goal(slide, "Goal: serve one quantized LLM across changing memory, latency, and quality budgets.")
    add_card(
        slide,
        1.18,
        2.55,
        3.05,
        2.70,
        "Deployment bottleneck",
        [
            "Autoregressive decode is dominated by repeated weight movement.",
            "Weight quantization cuts bandwidth and inference cost.",
            "At int2, quality often collapses.",
        ],
        accent=BLUE,
        fill=LIGHT_BLUE,
    )
    add_card(
        slide,
        4.45,
        2.55,
        3.05,
        2.70,
        "Past research direction",
        [
            "GPTQ, AWQ, SmoothQuant, OmniQuant, and QAT optimize one target bit-width.",
            "Each precision is usually a separate optimization problem.",
            "Multi-scale nesting had not been shown at LLM scale.",
        ],
        accent=ORANGE,
        fill=LIGHT_ORANGE,
    )
    add_card(
        slide,
        7.72,
        2.55,
        3.05,
        2.70,
        "What was missing",
        [
            "One checkpoint that can expose int8, int4, int2, and interpolated bit-widths.",
            "Low-bit slices that remain accurate.",
            "Dense accuracy-cost trade-offs without retraining.",
        ],
        accent=NAVY,
        fill=WHITE,
    )
    add_callout(slide, "MatQuant asks: can lower-precision models be nested inside one trained high-precision model?")
    add_source(slide)


def build_slide_2(prs: Presentation) -> None:
    slide = add_slide(prs, "Core Idea: Train Shared Most-Significant Bits", 2)
    add_goal(slide, "The key observation is structural: integer bit-widths already contain a natural nesting.")

    x0, y0 = 1.3, 3.0
    bit_colors = [NAVY, NAVY, NAVY, NAVY, BLUE, BLUE, ORANGE, ORANGE]
    for i in range(8):
        add_label(slide, str(7 - i), x0 + i * 0.47, y0, 0.40, 0.42, bit_colors[i], font_size=10)
    add_textbox(slide, "int8 parent weight", x0 + 0.35, y0 - 0.43, 2.7, 0.25, font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "slice MSBs", x0 + 4.0, y0 + 0.03, 1.2, 0.24, font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_arrow(slide, x0 + 3.95, y0 + 0.21, x0 + 4.95, y0 + 0.21, color=GRAY)

    for i in range(4):
        add_label(slide, str(7 - i), 6.55 + i * 0.47, 2.62, 0.40, 0.42, NAVY, font_size=10)
    add_textbox(slide, "int4 slice", 6.70, 2.18, 1.55, 0.25, font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    for i in range(2):
        add_label(slide, str(7 - i), 7.02 + i * 0.47, 3.58, 0.40, 0.42, ORANGE, font_size=10)
    add_textbox(slide, "int2 slice", 6.93, 4.06, 1.25, 0.25, font_size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    add_card(
        slide,
        1.25,
        4.45,
        2.85,
        1.75,
        "1. Parent model",
        ["Quantize at the largest precision, e.g., int8.", "Make early bit planes carry shared information."],
        accent=BLUE,
    )
    add_card(
        slide,
        4.33,
        4.45,
        2.85,
        1.75,
        "2. Slicing operator",
        ["Extract lower precision by right/left shifting MSBs.", "No separate checkpoint is required."],
        accent=ORANGE,
    )
    add_card(
        slide,
        7.42,
        4.45,
        2.85,
        1.75,
        "3. Joint loss",
        ["Optimize int8, int4, and int2 losses together.", "Works with OmniQuant or QAT parameters."],
        accent=GREEN,
    )
    add_callout(slide, "A smaller model is not approximated after training; it is co-trained as a slice of the parent.")
    add_source(slide)


def build_slide_3(prs: Presentation) -> None:
    slide = add_slide(prs, "Core Observations: What Makes Low Bits Work", 3)
    add_goal(slide, "MatQuant improves low-bit quality because training reshapes where information lives.")
    add_card(
        slide,
        1.18,
        2.48,
        2.62,
        1.75,
        "Shared MSBs matter",
        ["Slicing a normal int8 model is weak.", "The MSBs must be trained for low-bit use."],
        accent=NAVY,
        fill=WHITE,
    )
    add_card(
        slide,
        4.02,
        2.48,
        2.62,
        1.75,
        "Weight usage shifts",
        ["Histograms move toward larger quantized values.", "That gives int2 more useful buckets."],
        accent=BLUE,
        fill=LIGHT_BLUE,
    )
    add_card(
        slide,
        6.86,
        2.48,
        2.62,
        1.75,
        "Loss re-weighting",
        ["A higher int2 loss weight is essential.", "Optimizing int8/int4 too hard can hurt int2."],
        accent=ORANGE,
        fill=LIGHT_ORANGE,
    )
    add_card(
        slide,
        9.70,
        2.48,
        2.62,
        1.75,
        "Co-distillation",
        ["Higher-precision slices supervise lower-precision slices.", "int2 gains from an internal teacher."],
        accent=GREEN,
        fill=LIGHT_GREEN,
    )

    chart_x, chart_y = 2.08, 5.15
    add_textbox(slide, "Illustration: MatQuant raises low-bit accuracy while keeping high-bit quality", chart_x, chart_y - 0.40, 6.4, 0.28, font_size=12, bold=True, color=NAVY)
    values = [65, 67, 72, 74]
    labels = ["base", "MatQ", "int4", "int8"]
    colors = [RED, ORANGE, BLUE, NAVY]
    for i, (value, label, color) in enumerate(zip(values, labels, colors)):
        h = (value - 55) / 20 * 1.55
        x = chart_x + 0.35 + i * 1.15
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(chart_y + 1.70 - h), Inches(0.56), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        add_textbox(slide, f"{value}", x - 0.05, chart_y + 1.45 - h, 0.68, 0.22, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x - 0.20, chart_y + 1.82, 0.95, 0.45, font_size=8.2, color=BLACK, align=PP_ALIGN.CENTER)
    add_arrow(slide, chart_x + 0.25, chart_y + 1.72, chart_x + 5.0, chart_y + 1.72, color=GRAY, width=1.0)
    add_arrow(slide, chart_x + 0.25, chart_y + 1.72, chart_x + 0.25, chart_y + 0.05, color=GRAY, width=1.0)
    add_textbox(slide, "directional summary of paper tables", chart_x + 5.35, chart_y + 1.48, 2.9, 0.26, font_size=8.5, italic=True, color=GRAY)
    add_callout(slide, "The important effect is not compression alone; it is training the same bits to serve multiple precisions.")
    add_source(slide)


def build_slide_4(prs: Presentation) -> None:
    slide = add_slide(prs, "Main Results and Takeaway", 4)
    add_goal(slide, "MatQuant preserves high-bit accuracy and materially improves int2 slices.")
    add_card(
        slide,
        1.18,
        2.44,
        4.55,
        1.45,
        "OmniQuant results",
        [
            "int8 and int4 are within about 0.5% of separately trained baselines.",
            "int2 improves by 1.04%, 3.11%, and 3.01% on Gemma-2 2B, Gemma-2 9B, and Mistral 7B.",
        ],
        accent=BLUE,
        fill=LIGHT_BLUE,
        title_size=14,
        body_size=10.8,
    )
    add_card(
        slide,
        6.12,
        2.44,
        4.55,
        1.45,
        "QAT results",
        [
            "The same MatQuant recipe also works when weights are updated by QAT.",
            "int2 gains reach 4.46%, 6.27%, and 7.02% across the same model families.",
        ],
        accent=ORANGE,
        fill=LIGHT_ORANGE,
        title_size=14,
        body_size=10.8,
    )
    add_card(
        slide,
        1.18,
        4.28,
        4.55,
        1.45,
        "Elastic serving",
        [
            "int3 and int6 can be sliced even though they are not direct training targets.",
            "Layer-wise Mix'n'Match creates many deployment points from one model.",
        ],
        accent=GREEN,
        fill=LIGHT_GREEN,
        title_size=14,
        body_size=10.8,
    )
    add_card(
        slide,
        6.12,
        4.28,
        4.55,
        1.45,
        "Extra precision outlier bucket",
        [
            "A 2.05-bit variant adds a rare extra bucket for outliers.",
            "On Gemma-2 9B, the paper reports about 5% additional accuracy, boosted to 6% with co-distillation.",
        ],
        accent=PURPLE,
        fill=WHITE,
        title_size=14,
        body_size=10.4,
    )
    add_callout(slide, "Takeaway: one trained quantized model can expose a dense accuracy-cost frontier at deployment time.")
    add_source(slide)


def main() -> None:
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    for builder in (build_slide_1, build_slide_2, build_slide_3, build_slide_4):
        builder(prs)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
