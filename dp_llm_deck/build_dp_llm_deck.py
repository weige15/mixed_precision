from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "to_human"
OUT = OUT_DIR / "DP_LLM_4page_intro.pptx"
TEMPLATE = ROOT / "template.pptx"

NAVY = RGBColor(8, 56, 118)
BLUE = RGBColor(91, 155, 213)
LIGHT_BLUE = RGBColor(221, 236, 249)
ORANGE = RGBColor(237, 125, 49)
LIGHT_ORANGE = RGBColor(255, 235, 220)
YELLOW = RGBColor(243, 226, 127)
PALE_YELLOW = RGBColor(255, 246, 196)
GRAY = RGBColor(92, 92, 92)
LIGHT_GRAY = RGBColor(242, 244, 247)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(112, 173, 71)
LIGHT_GREEN = RGBColor(230, 244, 222)


def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def remove_shape(shape) -> None:
    el = shape._element  # pylint: disable=protected-access
    el.getparent().remove(el)


def disable_bullet(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()  # pylint: disable=protected-access
    for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buNone"):
        for child in list(p_pr.findall(qn(tag))):
            p_pr.remove(child)
    p_pr.append(OxmlElement("a:buNone"))


def add_slide(prs: Presentation, title: str, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[17])
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx != 0:
            remove_shape(shape)
    title_shape = slide.shapes.title
    title_shape.text = title
    tf = title_shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    for p in tf.paragraphs:
        disable_bullet(p)
        p.font.name = "Arial"
        p.font.size = Pt(29)
        p.font.color.rgb = BLACK
    add_textbox(slide, str(slide_no), 15.42, 8.28, 0.42, 0.30, font_size=10, color=WHITE)
    return slide


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 14,
    color=BLACK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return box


def set_runs(paragraph, font_size=14, color=BLACK, bold=False):
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    for run in paragraph.runs:
        if run.font is None:
            continue
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_goal(slide, text: str, y: float = 1.72):
    add_textbox(slide, text, 1.12, y, 9.2, 0.35, font_size=15, color=NAVY, bold=True)
    start_x = 1.12
    for i in range(8):
        dash = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(start_x + i * 0.82),
            Inches(y + 0.43),
            Inches(0.58),
            Inches(0.07),
        )
        dash.fill.solid()
        dash.fill.fore_color.rgb = BLUE
        dash.line.fill.background()


def add_card(slide, x, y, w, h, title, bullets, accent=NAVY, fill=WHITE, title_size=15, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, title_size, accent, True)
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.space_before = Pt(5)
        p.level = 0
        set_runs(p, body_size, BLACK, False)
    return shape


def add_callout(slide, text, x=2.2, y=7.08, w=7.6, h=0.45):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = YELLOW
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, 13, NAVY, True)
    return shape


def add_source(slide, text):
    add_textbox(slide, text, 1.12, 8.02, 9.4, 0.35, font_size=10, color=GRAY)


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.4):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_bit_grid(slide, x, y, values, colors, cell_w=0.25, cell_h=0.22):
    for r, row in enumerate(values):
        for c, val in enumerate(row):
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(x + c * cell_w),
                Inches(y + r * cell_h),
                Inches(cell_w - 0.01),
                Inches(cell_h - 0.01),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = colors.get(val, LIGHT_GRAY)
            rect.line.color.rgb = WHITE


def build_slide_1(prs):
    slide = add_slide(prs, "DP-LLM: Problem Setting and Research Gap", 1)
    add_goal(slide, "Goal: adapt one on-device LLM to changing latency, memory, and quality constraints.")
    add_card(
        slide,
        1.18,
        2.55,
        3.05,
        2.70,
        "Past research: static PTQ",
        [
            "Weight-only PTQ cuts memory traffic.",
            "Precision is fixed after calibration.",
            "Cannot react to per-query runtime budgets.",
        ],
        accent=BLUE,
        fill=WHITE,
        body_size=12,
    )
    add_card(
        slide,
        4.48,
        2.55,
        3.05,
        2.70,
        "Past research: multi-scale",
        [
            "Multiple bit-width views share storage.",
            "Layer-wise mixed precision supports non-integer bits.",
            "Assignments remain offline or static during decoding.",
        ],
        accent=ORANGE,
        fill=WHITE,
        body_size=12,
    )
    add_card(
        slide,
        7.78,
        2.55,
        3.05,
        2.70,
        "What is missing",
        [
            "How to match target precision or latency at runtime.",
            "Layer sensitivity changes token by token.",
            "Static precision leaves trade-offs unused.",
        ],
        accent=NAVY,
        fill=WHITE,
        body_size=12,
    )
    add_callout(slide, "DP-LLM treats precision as a runtime, layer-wise decision")
    add_source(slide, "Source: Kwon et al., DP-LLM, NeurIPS 2025 / arXiv:2508.06041")


def build_slide_2(prs):
    slide = add_slide(prs, "Core Observation: Sensitivity Changes by Token", 2)
    add_goal(slide, "The high-precision layers are not fixed; sensitivity shifts across decoding steps.")
    labels = [
        ("Uniform", "same bit-width", 1.25, BLUE),
        ("Static layer-wise", "fixed by layer", 4.45, ORANGE),
        ("Dynamic layer-wise", "changes by token", 7.65, NAVY),
    ]
    grids = [
        [[3, 3, 3, 3, 3, 3], [3, 3, 3, 3, 3, 3], [3, 3, 3, 3, 3, 3], [3, 3, 3, 3, 3, 3], [3, 3, 3, 3, 3, 3]],
        [[3, 3, 3, 3, 3, 3], [4, 4, 4, 4, 4, 4], [3, 3, 3, 3, 3, 3], [5, 5, 5, 5, 5, 5], [4, 4, 4, 4, 4, 4]],
        [[3, 4, 3, 5, 3, 4], [4, 3, 5, 3, 4, 3], [3, 5, 4, 3, 3, 5], [5, 3, 3, 4, 5, 3], [4, 5, 3, 4, 3, 5]],
    ]
    color_map = {3: LIGHT_BLUE, 4: PALE_YELLOW, 5: LIGHT_ORANGE}
    for idx, (name, desc, x, accent) in enumerate(labels):
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.48),
            Inches(2.65),
            Inches(2.65),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = accent
        panel.line.width = Pt(1.5)
        add_textbox(slide, name, x + 0.20, 2.70, 2.25, 0.32, font_size=14, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, x + 0.28, 3.08, 2.08, 0.46, font_size=10, color=BLACK, align=PP_ALIGN.CENTER)
        add_textbox(slide, "Layer", x + 0.16, 3.63, 0.70, 0.28, font_size=10, color=GRAY)
        add_bit_grid(slide, x + 0.60, 3.92, grids[idx], color_map, cell_h=0.18)
        add_textbox(slide, "Step ->", x + 0.96, 4.84, 1.25, 0.28, font_size=10, color=GRAY)
    add_arrow(slide, 3.98, 3.80, 4.25, 3.80, color=GRAY)
    add_arrow(slide, 7.20, 3.80, 7.45, 3.80, color=GRAY)
    add_callout(slide, "Key observation: dynamic layer-wise mixed precision captures token-step sensitivity", x=1.72, y=6.58, w=8.35, h=0.56)
    add_textbox(
        slide,
        "The paper uses an oracle analysis to show that dynamic assignment can lower perplexity versus static layer-wise mixed precision.",
        1.30,
        7.22,
        9.2,
        0.38,
        font_size=11,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Source: paper Figures 2-3; Section 2.4")


def add_phase(slide, x, y, title, bullets, accent):
    shape = add_card(slide, x, y, 2.30, 1.78, title, bullets, accent=accent, fill=WHITE, title_size=12, body_size=10)
    return shape


def build_slide_3(prs):
    slide = add_slide(prs, "Method: DP-LLM Precision Selector", 3)
    add_goal(slide, "Offline thresholds plus runtime selectors choose low or high precision per layer.")
    phases = [
        ("1. Fit budget", ["Choose each layer's max precision.", "Use static sensitivity under memory budget."], BLUE),
        ("2. Assign p_i", ["Fine-tune average precision scalar.", "Regularize model average to target bits."], ORANGE),
        ("3. Threshold T_i", ["Use calibration relative-error distribution.", "Map p_i to low/high precision rate."], GREEN),
        ("4. Runtime select", ["Estimate relative error cheaply.", "Compare with T_i, then choose W_l or W_h."], NAVY),
    ]
    xs = [1.05, 3.58, 6.10, 8.62]
    for i, (title, bullets, color) in enumerate(phases):
        add_phase(slide, xs[i], 2.30, title, bullets, color)
        if i < 3:
            add_arrow(slide, xs[i] + 2.32, 3.20, xs[i + 1] - 0.10, 3.20, color=GRAY)
    add_textbox(slide, "Runtime selector inside each linear layer", 1.14, 4.68, 4.0, 0.28, font_size=13, color=NAVY, bold=True)
    add_textbox(slide, "x", 1.28, 5.55, 0.35, 0.28, font_size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 1.68, 5.68, 2.20, 5.68)
    add_card(slide, 2.20, 5.15, 1.75, 1.05, "Relative error estimator", ["approx. ||(W_h-W_l)x||"], accent=BLUE, fill=LIGHT_BLUE, title_size=10, body_size=10)
    add_arrow(slide, 3.96, 5.68, 4.45, 5.68)
    add_card(slide, 4.45, 5.15, 1.55, 1.05, "Compare", ["error > T_i ?"], accent=ORANGE, fill=LIGHT_ORANGE, title_size=10, body_size=10)
    add_arrow(slide, 6.00, 5.68, 6.55, 5.40, color=GREEN)
    add_arrow(slide, 6.00, 5.68, 6.55, 5.98, color=NAVY)
    add_card(slide, 6.55, 4.86, 1.45, 0.70, "Low bit", ["W_l"], accent=GREEN, fill=LIGHT_GREEN, title_size=10, body_size=10)
    add_card(slide, 6.55, 5.80, 1.45, 0.70, "High bit", ["W_h"], accent=NAVY, fill=LIGHT_BLUE, title_size=10, body_size=10)
    add_arrow(slide, 8.02, 5.20, 8.72, 5.45, color=GRAY)
    add_arrow(slide, 8.02, 6.15, 8.72, 5.92, color=GRAY)
    add_card(slide, 8.70, 5.25, 1.75, 0.95, "GEMV output", ["selected precision"], accent=NAVY, fill=WHITE, title_size=10, body_size=10)
    add_textbox(
        slide,
        "Efficiency trick: use linear-regression estimation when ||x|| predicts error; otherwise use random projection. Async estimation hides part of the cost.",
        1.12,
        7.15,
        9.55,
        0.66,
        font_size=10,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Source: paper Sections 3-5; Figures 4-6")


def build_slide_4(prs):
    slide = add_slide(prs, "Main Results and Takeaways", 4)
    add_goal(slide, "DP-LLM improves quality versus static mixed precision with small runtime overhead.")
    add_card(
        slide,
        1.12,
        2.42,
        3.10,
        2.82,
        "Quality",
        [
            "Models: Llama-3-8B and Phi-3-Medium.",
            "Lower PPL than LLM-MQ / HAWQ-V2 on WikiText2 and C4 in most settings.",
            "Often strongest on GSM8K, MBPP, BBH, and MATH.",
        ],
        accent=BLUE,
        fill=WHITE,
        body_size=11,
    )
    add_card(
        slide,
        4.45,
        2.42,
        3.10,
        2.82,
        "Runtime overhead",
        [
            "Selector overhead geomean: 1.45% for Llama-3-8B; 0.81% for Phi-3-Medium.",
            "RTX 4060Ti Hybrid+Async: 0.74%, 0.66%, 0.45% at 3.5/4.0/4.5 bits.",
        ],
        accent=ORANGE,
        fill=WHITE,
        body_size=11,
    )
    add_card(
        slide,
        7.78,
        2.42,
        3.10,
        2.82,
        "Conclusion",
        [
            "Relative error works as a precision-choice proxy.",
            "Dynamic assignment captures sensitivity missed by static baselines.",
            "99th percentile effective-bit drift is about 2.25-3.32%.",
        ],
        accent=NAVY,
        fill=WHITE,
        body_size=11,
    )
    add_callout(slide, "Takeaway: target precision becomes low-overhead per-layer runtime adaptation", x=1.80, y=6.72, w=8.20, h=0.56)
    add_textbox(slide, "Benchmarks: WikiText2, C4, GSM8K, MBPP, BBH, MATH; hardware: Jetson Orin AGX and RTX 4060 Ti.", 1.22, 7.42, 9.35, 0.30, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: paper Tables 1-7; Section 6")


def main() -> None:
    prs = Presentation(TEMPLATE)
    clear_slides(prs)
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    OUT_DIR.mkdir(exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
