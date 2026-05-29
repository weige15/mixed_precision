from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
BASE = ROOT / "mobiquant_base.pptx"
OUT = ROOT / "mobiquant_summary.pptx"

NAVY = RGBColor(9, 54, 112)
BLUE = RGBColor(91, 155, 213)
LIGHT_BLUE = RGBColor(226, 239, 255)
YELLOW = RGBColor(242, 224, 132)
GREEN = RGBColor(47, 125, 50)
MAGENTA = RGBColor(205, 96, 201)
RED = RGBColor(216, 72, 98)
ORANGE = RGBColor(242, 122, 46)
GRAY = RGBColor(245, 247, 250)
DARK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def clear_generated_body(slide):
    for shape in list(slide.shapes):
        if shape.name != "DeckTitle":
            delete_shape(shape)


def set_title(slide, text):
    title = slide.shapes.title
    title.name = "DeckTitle"
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p_pr = p._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in {qn("a:buNone"), qn("a:buAutoNum"), qn("a:buChar"), qn("a:buBlip")}:
            p_pr.remove(child)
    p_pr.insert(0, OxmlElement("a:buNone"))
    p.font.name = "Arial"
    p.font.size = Pt(31)
    p.font.bold = False
    p.font.color.rgb = DARK


def add_textbox(slide, x, y, w, h, text="", font_size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_page_number(slide, num):
    add_textbox(slide, 15.42, 8.28, 0.42, 0.3, str(num), font_size=10, bold=False, color=NAVY, align=PP_ALIGN.CENTER)


def add_card(slide, x, y, w, h, title, bullets, accent=NAVY):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.4)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = accent
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(12.6)
        p.font.color.rgb = DARK
        p.space_before = Pt(4)
    return card


def add_label(slide, x, y, w, h, text, fill, font_size=12, color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.6):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_source(slide, text="Source: Wang et al., MoBiQuant, arXiv:2602.20191 (2026)"):
    add_textbox(slide, 1.18, 8.18, 9.4, 0.26, text, font_size=8.5, color=RGBColor(90, 90, 90))


def slide_problem(slide):
    set_title(slide, "MoBiQuant: Problem Setting and Research Gap")
    clear_generated_body(slide)
    add_textbox(
        slide, 1.18, 2.02, 9.4, 0.45,
        "Goal: one quantized LLM should adapt precision online as latency and memory budgets change.",
        font_size=17, bold=True, color=NAVY,
    )
    add_card(slide, 1.18, 2.72, 3.05, 3.0, "Past research: static PTQ", [
        "GPTQ, OmniQuant, SmoothQuant, AWQ tune one target bit-width.",
        "They optimize fixed-precision quality, not runtime elasticity.",
        "Changing bit-width usually requires recalibration or quality loss.",
    ], accent=BLUE)
    add_card(slide, 4.45, 2.72, 3.05, 3.0, "Past research: any-precision", [
        "AnyPrecisionLLM, AnyBCQ, MatQuant support multiple precisions.",
        "But they add scaling, repacking, table lookups, or kernel overhead.",
        "Switching precision is not yet cheap at token granularity.",
    ], accent=ORANGE)
    add_card(slide, 7.72, 2.72, 3.05, 3.0, "What is missing", [
        "A single nested model representation.",
        "Token-level precision decisions during inference.",
        "A low-overhead kernel path that fetches only needed bits.",
    ], accent=NAVY)
    add_label(slide, 2.45, 6.35, 7.1, 0.48, "MoBiQuant frames precision as a token-level runtime resource", YELLOW, font_size=14, color=NAVY)
    add_source(slide)


def slide_observation(slide):
    set_title(slide, "Core Observation: Outlier Migration")
    clear_generated_body(slide)
    add_textbox(
        slide, 1.18, 2.0, 9.35, 0.42,
        "The tokens that dominate quantization error move when the bit-width changes.",
        font_size=17, bold=True, color=NAVY,
    )
    add_card(slide, 1.18, 2.65, 4.6, 1.35, "Why static calibration fails", [
        "Calibration parameters overfit the outlier distribution of one bit-width.",
        "A token well-fitted at 4-bit can become an outlier at 3-bit.",
    ], accent=RED)
    add_card(slide, 6.02, 2.65, 4.65, 1.35, "What MoBiQuant exploits", [
        "Token sensitivity can guide per-token bit allocation.",
        "Routing high-sensitivity tokens through more bit slices improves generalization.",
    ], accent=GREEN)

    chart_x, chart_y = 1.35, 4.65
    add_textbox(slide, chart_x, chart_y - 0.32, 4.2, 0.25, "LLaMA3-8B / WikiText2 PPL", font_size=10.5, bold=True, color=DARK)
    values = [10.06, 9.01, 7.41, 7.31]
    labels = ["3b->4b", "retain 10%", "4b->4b", "MoBiQ"]
    colors = [GREEN, MAGENTA, RED, ORANGE]
    max_v = 10.6
    base_y = chart_y + 2.55
    for i, (v, lab, col) in enumerate(zip(values, labels, colors)):
        x = chart_x + i * 0.88
        h = 2.25 * v / max_v
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(base_y - h), Inches(0.5), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.color.rgb = col
        add_textbox(slide, x - 0.06, base_y - h + 0.1, 0.62, 0.25, f"{v:.2f}", font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - 0.22, base_y + 0.08, 1.0, 0.3, lab, font_size=7.0, color=DARK, align=PP_ALIGN.CENTER)
    add_arrow(slide, chart_x - 0.05, base_y, chart_x + 3.6, base_y, color=DARK, width=1.0)
    add_arrow(slide, chart_x - 0.05, base_y, chart_x - 0.05, chart_y + 0.15, color=DARK, width=1.0)
    add_textbox(slide, chart_x + 0.86, chart_y + 0.55, 0.95, 0.32, "2.65 PPL gap", font_size=8.0, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_arrow(slide, chart_x + 1.58, chart_y + 0.75, chart_x + 2.25, chart_y + 1.45, color=RED, width=1.4)

    plot_x, plot_y = 6.25, 4.55
    add_textbox(slide, plot_x, plot_y - 0.22, 4.3, 0.25, "Token-wise error peaks shift across precision", font_size=10.5, bold=True, color=DARK)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(plot_x), Inches(plot_y + 0.15), Inches(4.15), Inches(2.25))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(160, 160, 160)
    for idx in range(18):
        x = plot_x + 0.18 + idx * 0.21
        h1 = [0.35, 0.7, 0.45, 1.55, 0.4, 0.28, 0.9, 0.32, 0.5, 0.42, 1.85, 0.36, 0.52, 0.4, 1.25, 0.31, 0.48, 0.38][idx]
        h2 = [0.4, 0.32, 1.65, 0.38, 0.7, 0.35, 0.3, 1.95, 0.36, 0.48, 0.35, 1.42, 0.33, 0.75, 0.44, 1.7, 0.34, 0.42][idx]
        s1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(plot_y + 2.28 - h1), Inches(0.06), Inches(h1))
        s1.fill.solid(); s1.fill.fore_color.rgb = BLUE; s1.line.color.rgb = BLUE
        s2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.08), Inches(plot_y + 2.28 - h2), Inches(0.06), Inches(h2))
        s2.fill.solid(); s2.fill.fore_color.rgb = GREEN; s2.line.color.rgb = GREEN
    add_label(slide, plot_x + 0.3, plot_y + 2.55, 1.05, 0.28, "4-bit errors", BLUE, font_size=8)
    add_label(slide, plot_x + 1.5, plot_y + 2.55, 1.05, 0.28, "3-bit errors", GREEN, font_size=8)
    add_label(slide, plot_x + 2.88, plot_y + 2.55, 1.0, 0.28, "low overlap", RED, font_size=8)
    add_source(slide)


def slide_method(slide):
    set_title(slide, "Method: MoBiSlice and MoBiRoute")
    clear_generated_body(slide)
    add_textbox(slide, 1.18, 2.02, 9.5, 0.38, "A single linear block becomes a routed mixture of residual bit slices.", font_size=17, bold=True, color=NAVY)

    add_label(slide, 1.28, 3.0, 1.5, 0.5, "Input tokens", BLUE, font_size=12)
    add_label(slide, 3.25, 3.0, 1.55, 0.5, "MoBiRoute MLP", ORANGE, font_size=10)
    add_label(slide, 5.35, 3.0, 1.95, 0.5, "Binary slice mask", RED, font_size=11)
    add_label(slide, 7.9, 3.0, 2.25, 0.5, "Active bit slices", GREEN, font_size=11)
    add_arrow(slide, 2.78, 3.25, 3.22, 3.25)
    add_arrow(slide, 4.8, 3.25, 5.32, 3.25)
    add_arrow(slide, 7.3, 3.25, 7.87, 3.25)

    for i, tok in enumerate(["x1", "x2", "x3"]):
        add_label(slide, 1.38 + i * 0.38, 3.75, 0.3, 0.3, tok, LIGHT_BLUE, font_size=8, color=NAVY)
    for row in range(3):
        for col in range(4):
            fill = RED if (row, col) in [(0, 0), (0, 1), (1, 0), (1, 1), (1, 2), (2, 0)] else WHITE
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(5.55 + col * 0.35), Inches(3.72 + row * 0.28), Inches(0.3), Inches(0.23))
            cell.fill.solid(); cell.fill.fore_color.rgb = fill; cell.line.color.rgb = RGBColor(160, 160, 160)
    for i, lab in enumerate(["W1", "W2", "W3", "W4"]):
        y = 3.78 + i * 0.38
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.05 + i * 0.34), Inches(y), Inches(1.0), Inches(0.34))
        s.fill.solid()
        s.fill.fore_color.rgb = [BLUE, ORANGE, GREEN, RGBColor(220, 220, 220)][i]
        s.line.color.rgb = NAVY
        add_textbox(slide, 8.08 + i * 0.34, y + 0.04, 0.9, 0.16, lab, font_size=8.5, bold=True, color=WHITE if i < 3 else DARK, align=PP_ALIGN.CENTER)

    add_card(slide, 1.18, 5.25, 3.05, 1.65, "MoBiSlice", [
        "Recursively quantizes residual weights into 2-bit slices.",
        "Higher precision is reconstructed by summing more slices.",
    ], accent=BLUE)
    add_card(slide, 4.45, 5.25, 3.05, 1.65, "MoBiRoute", [
        "Scores token sensitivity and gates slices with binary decisions.",
        "A threshold controls the target average bit budget.",
    ], accent=ORANGE)
    add_card(slide, 7.72, 5.25, 3.05, 1.65, "Kernel design", [
        "Bit-major packing fetches only active slices.",
        "Shared scaling avoids extra per-precision parameters.",
    ], accent=GREEN)
    add_source(slide)


def slide_results(slide):
    set_title(slide, "Main Results and Takeaways")
    clear_generated_body(slide)
    add_textbox(slide, 1.18, 2.02, 9.45, 0.38, "MoBiQuant keeps elasticity while approaching fixed-precision PTQ quality.", font_size=17, bold=True, color=NAVY)
    add_card(slide, 1.18, 2.75, 4.55, 1.45, "Accuracy", [
        "Matches or surpasses static scalar PTQ across LLaMA2 and LLaMA3 model families.",
        "At 2-3 bit, reduces severe PPL collapse versus prior any-precision baselines.",
    ], accent=BLUE)
    add_card(slide, 6.12, 2.75, 4.55, 1.45, "Throughput", [
        "Average decoding speedup: 33.8% over AnyPrecisionLLM.",
        "Average decoding speedup: 22.8% over AnyBCQ.",
    ], accent=ORANGE)
    add_card(slide, 1.18, 4.55, 4.55, 1.45, "Memory", [
        "One nested model replaces multiple per-bit deployments.",
        "Reported memory footprint reduction is up to 3.5x.",
    ], accent=GREEN)
    add_card(slide, 6.12, 4.55, 4.55, 1.45, "Conclusion", [
        "The core win is not just storing more precisions.",
        "The router mitigates outlier migration by assigning precision per token.",
    ], accent=NAVY)
    add_label(slide, 2.25, 6.55, 7.3, 0.55, "Takeaway: precision becomes a dynamic information budget at inference time", YELLOW, font_size=14, color=NAVY)
    add_source(slide)


def main():
    prs = Presentation(str(BASE))
    builders = [slide_problem, slide_observation, slide_method, slide_results]
    for idx, (slide, builder) in enumerate(zip(prs.slides, builders), start=1):
        builder(slide)
        add_page_number(slide, idx)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
