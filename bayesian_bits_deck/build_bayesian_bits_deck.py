from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PROJECT = ROOT.parent
TEMPLATE = PROJECT / "template.pptx"
OUT = ROOT / "bayesian_bits_4page_intro.pptx"
OUT_HUMAN = PROJECT / "to_human" / "Bayesian_Bits_4page_intro.pptx"

NAVY = RGBColor(9, 54, 112)
BLUE = RGBColor(65, 126, 211)
LIGHT_BLUE = RGBColor(229, 241, 255)
YELLOW = RGBColor(242, 224, 132)
GREEN = RGBColor(48, 132, 82)
MAGENTA = RGBColor(182, 86, 184)
RED = RGBColor(210, 76, 94)
ORANGE = RGBColor(232, 126, 47)
GRAY = RGBColor(245, 247, 250)
DARK = RGBColor(0, 0, 0)
MID = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)


def delete_slide(prs, slide):
    slide_id = slide.slide_id
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sld_id in slides:
        if int(sld_id.get("id")) == slide_id:
            prs.part.drop_rel(sld_id.get(qn("r:id")))
            xml_slides.remove(sld_id)
            return


def new_deck():
    prs = Presentation(str(TEMPLATE))
    for slide in list(prs.slides):
        delete_slide(prs, slide)
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[17])
    return prs


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def clear_generated_body(slide):
    for shape in list(slide.shapes):
        if shape.name != "DeckTitle":
            delete_shape(shape)


def set_title(slide, text):
    title = slide.shapes.title
    title.name = "DeckTitle"
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p_pr = p._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in {qn("a:buNone"), qn("a:buAutoNum"), qn("a:buChar"), qn("a:buBlip")}:
            p_pr.remove(child)
    p_pr.insert(0, OxmlElement("a:buNone"))
    p.font.name = "Arial"
    p.font.size = Pt(31)
    p.font.bold = False
    p.font.color.rgb = DARK


def add_textbox(slide, x, y, w, h, text="", font_size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_page_number(slide, num):
    add_textbox(slide, 15.42, 8.28, 0.42, 0.3, str(num), font_size=10, color=NAVY, align=PP_ALIGN.CENTER)


def add_source(slide):
    add_textbox(
        slide, 1.18, 8.18, 10.0, 0.26,
        "Source: van Baalen et al., Bayesian Bits: Unifying Quantization and Pruning, NeurIPS 2020 / arXiv:2005.07093",
        font_size=8.3, color=MID,
    )


def add_card(slide, x, y, w, h, title, bullets, accent=NAVY, title_size=15, body_size=11.4):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.35)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.17)
    tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.11)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = accent
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.name = "Arial"
        p.font.size = Pt(body_size)
        p.font.color.rgb = DARK
        p.space_before = Pt(4)
    return card


def add_label(slide, x, y, w, h, text, fill, font_size=12, color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def slide_problem(slide):
    set_title(slide, "Bayesian Bits: Problem and Research Gap")
    clear_generated_body(slide)
    add_textbox(
        slide, 1.18, 2.02, 9.75, 0.45,
        "Goal: learn an accuracy-efficiency policy instead of hand-picking pruning and bit-width settings.",
        font_size=17, bold=True, color=NAVY,
    )
    add_card(slide, 1.18, 2.75, 3.05, 2.85, "Deployment bottleneck", [
        "Inference cost depends on arithmetic and data movement.",
        "Quantization lowers operand width; pruning removes operations.",
        "The useful target is a Pareto trade-off, not one compressed model.",
    ], accent=BLUE)
    add_card(slide, 4.45, 2.75, 3.05, 2.85, "Past direction", [
        "Fixed-bit QAT optimizes one selected precision.",
        "Mixed precision searches layer or tensor bit-widths.",
        "Pruning is usually treated as a separate compression task.",
    ], accent=ORANGE)
    add_card(slide, 7.72, 2.75, 3.05, 2.85, "Missing capability", [
        "The bit-width search space is exponential.",
        "Continuous learned bits may not map to power-of-two hardware.",
        "Pruning and quantization need one trainable formulation.",
    ], accent=NAVY)
    add_label(slide, 2.25, 6.28, 7.6, 0.52, "Bayesian Bits turns bit allocation and pruning into gated optimization", YELLOW, font_size=14, color=NAVY)
    add_source(slide)


def slide_prior(slide):
    set_title(slide, "Why Existing Mixed Precision Was Not Enough")
    clear_generated_body(slide)
    add_textbox(slide, 1.18, 2.02, 9.5, 0.42, "Earlier methods improved quantization, but each left a deployment gap.", font_size=17, bold=True, color=NAVY)
    add_card(slide, 1.18, 2.72, 2.32, 3.15, "Fixed-bit QAT", [
        "PACT, LSQ, and TQT learn clipping or step sizes.",
        "Strong at one precision.",
        "They do not decide where each tensor should spend bits.",
    ], accent=BLUE, title_size=14, body_size=10.7)
    add_card(slide, 3.72, 2.72, 2.32, 3.15, "Sensitivity search", [
        "HAWQ-style methods use Hessian signals to rank layers.",
        "They expose sensitivity.",
        "A discrete assignment step is still needed.",
    ], accent=GREEN, title_size=14, body_size=10.7)
    add_card(slide, 6.26, 2.72, 2.32, 3.15, "NAS / RL search", [
        "Search methods can include hardware feedback.",
        "They optimize bit policies.",
        "Cost and policy complexity grow with model size.",
    ], accent=ORANGE, title_size=14, body_size=10.7)
    add_card(slide, 8.80, 2.72, 2.32, 3.15, "Differentiable bits", [
        "DQ learns continuous bit-widths.",
        "Real hardware needs rounded widths.",
        "Rounding can erase expected savings.",
    ], accent=RED, title_size=14, body_size=10.7)
    add_label(slide, 2.15, 6.45, 7.9, 0.5, "The paper targets hardware-friendly power-of-two widths by construction", LIGHT_BLUE, font_size=13.5, color=NAVY)
    add_source(slide)


def slide_method(slide):
    set_title(slide, "Core Idea: Gated Residual Quantization")
    clear_generated_body(slide)
    add_textbox(slide, 1.18, 1.98, 9.6, 0.42, "A tensor is represented as a low-bit value plus optional quantized residuals.", font_size=17, bold=True, color=NAVY)
    y = 3.05
    add_label(slide, 1.18, y, 1.3, 0.46, "2-bit base", BLUE, font_size=11.5)
    add_arrow(slide, 2.5, y + 0.23, 3.05, y + 0.23)
    add_label(slide, 3.08, y, 1.45, 0.46, "+ residual 4", ORANGE, font_size=11.5)
    add_arrow(slide, 4.55, y + 0.23, 5.1, y + 0.23)
    add_label(slide, 5.13, y, 1.45, 0.46, "+ residual 8", GREEN, font_size=11.5)
    add_arrow(slide, 6.6, y + 0.23, 7.15, y + 0.23)
    add_label(slide, 7.18, y, 1.65, 0.46, "+ residual 16", MAGENTA, font_size=11.5)
    add_arrow(slide, 8.85, y + 0.23, 9.4, y + 0.23)
    add_label(slide, 9.43, y, 1.45, 0.46, "+ residual 32", NAVY, font_size=11.5)

    add_textbox(slide, 1.55, 3.78, 8.9, 0.38, "xq = z2 · (x2 + z4 · (eps4 + z8 · (eps8 + z16 · (eps16 + z32 · eps32))))", font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_card(slide, 1.18, 4.62, 2.32, 1.72, "Start low", [
        "Begin with a 2-bit grid.",
        "This is the lowest nonzero representation.",
    ], accent=BLUE, title_size=14, body_size=10.5)
    add_card(slide, 3.72, 4.62, 2.32, 1.72, "Add residuals", [
        "Each residual halves the quantization error scale.",
        "Effective precision doubles.",
    ], accent=ORANGE, title_size=14, body_size=10.5)
    add_card(slide, 6.26, 4.62, 2.32, 1.72, "Learn gates", [
        "Stochastic binary gates choose residuals.",
        "Hard-concrete relaxation enables gradients.",
    ], accent=GREEN, title_size=14, body_size=10.5)
    add_card(slide, 8.80, 4.62, 2.32, 1.72, "Unify pruning", [
        "A 0-bit gate sits before the 2-bit value.",
        "Turning it off prunes the tensor group.",
    ], accent=RED, title_size=14, body_size=10.5)
    add_label(slide, 2.05, 6.85, 8.0, 0.5, "Bayesian prior: penalize active gates, optionally in proportion to BOP cost", YELLOW, font_size=13.5, color=NAVY)
    add_source(slide)


def add_table(slide, x, y, w, h, rows):
    n = len(rows)
    row_h = h / n
    for i, (left, right, accent) in enumerate(rows):
        yy = y + i * row_h
        bg = GRAY if i % 2 == 0 else WHITE
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(yy), Inches(w), Inches(row_h - 0.03))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(221, 226, 232)
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(yy), Inches(0.13), Inches(row_h - 0.03))
        tag.fill.solid()
        tag.fill.fore_color.rgb = accent
        tag.line.color.rgb = accent
        add_textbox(slide, x + 0.25, yy + 0.08, 3.25, row_h - 0.1, left, font_size=11.4, bold=True, color=accent)
        add_textbox(slide, x + 3.55, yy + 0.08, w - 3.75, row_h - 0.1, right, font_size=10.6, color=DARK)


def slide_results(slide):
    set_title(slide, "Observations, Results, and Takeaway")
    clear_generated_body(slide)
    add_textbox(slide, 1.18, 1.98, 9.55, 0.42, "The learned gates behave like a resource allocation policy over precision and channels.", font_size=17, bold=True, color=NAVY)
    rows = [
        ("Low-bit priors apply useful pressure", "The variational objective regularizes active gates; the prior can be scaled by BOP contribution.", BLUE),
        ("Joint pruning + quantization wins", "ImageNet ResNet18 ablations show the combined curve beats pruning-only and quantization-only variants.", GREEN),
        ("Known heuristics emerge", "Aggressive settings push most tensors to 2-bit while often preserving first and last layers at higher precision.", ORANGE),
        ("Post-training use is possible", "On pretrained ResNet18, the method can learn gates, or gates plus scales, without updating weights.", MAGENTA),
        ("Cost remains a limitation", "One ResNet18 run used 30 Bayesian Bits epochs plus 10 fixed-gate fine-tuning epochs, about 70 hours on one Tesla V100.", RED),
    ]
    add_table(slide, 1.18, 2.68, 9.95, 4.35, rows)
    add_label(slide, 2.0, 7.32, 8.3, 0.5, "Takeaway: pruning becomes 0-bit quantization; precision becomes a learnable budget", YELLOW, font_size=13.5, color=NAVY)
    add_source(slide)


def main():
    prs = new_deck()
    builders = [slide_problem, slide_prior, slide_method, slide_results]
    for idx, (slide, builder) in enumerate(zip(prs.slides, builders), start=1):
        builder(slide)
        add_page_number(slide, idx)
    prs.save(str(OUT))
    OUT_HUMAN.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_HUMAN))
    print(OUT)
    print(OUT_HUMAN)


if __name__ == "__main__":
    main()
